#!/usr/bin/env python3
"""
EM2 KIT BUILDER — builds a lesson deck from a Resource Lab saved kit.

Unlike build_module.py (which parses four source .pptx files), this takes a flat
ordered list of kit items and treats POSITION as the lesson index: item 1 becomes
Lesson 1, item 2 becomes Lesson 2, and so on. Duplicates are preserved.

A kit holds ONE resource type, so slides for the other types render their normal
dropzone placeholders. All slide-drawing helpers are reused from build_module.
"""
import os, re, glob, shutil, subprocess, tempfile
import requests
from pptx import Presentation
from pptx.util import Emu

import build_module as BM


# ---------- fetching ----------
DRIVE_FILE_RE = re.compile(r'/file/d/([A-Za-z0-9_-]+)')
DRIVE_ID_RE   = re.compile(r'[?&]id=([A-Za-z0-9_-]+)')
LH3_RE        = re.compile(r'lh3\.googleusercontent\.com/d/([A-Za-z0-9_-]+)')


def _drive_id(url):
    """Pull a Drive file id out of any of the URL shapes we store."""
    for rx in (DRIVE_FILE_RE, LH3_RE, DRIVE_ID_RE):
        m = rx.search(url or "")
        if m:
            return m.group(1)
    return None


def _download(url, dest):
    """Fetch a kit item's bytes. Drive /preview and lh3 links are rewritten to
    their direct-download form; anything else is fetched as-is."""
    fid = _drive_id(url)
    if fid:
        url = f"https://drive.google.com/uc?export=download&id={fid}"
    r = requests.get(url, allow_redirects=True, timeout=120)
    r.raise_for_status()
    with open(dest, "wb") as f:
        f.write(r.content)
    return dest


def _pdf_pages_to_pngs(pdf, outdir, prefix, dpi=150):
    os.makedirs(outdir, exist_ok=True)
    subprocess.run(["pdftoppm", "-png", "-r", str(dpi), pdf,
                    os.path.join(outdir, prefix)], check=True)
    return sorted(glob.glob(os.path.join(outdir, prefix + "*.png")))


def _is_pdf(path):
    try:
        with open(path, "rb") as f:
            return f.read(5) == b"%PDF-"
    except Exception:
        return False


# ---------- per-item asset prep ----------
def _dare_text(pdf):
    """Pull the Question and Words text off page 1 of a DARE PDF.

    These fields are not stored in lab_items -- they only exist inside the PDF --
    but the decks were generated from PPTX so page 1 has a real text layer and
    needs no OCR. Returns ("", "") if anything is missing rather than raising:
    the slide builder renders an em dash for empty strings.
    """
    try:
        out = subprocess.run(["pdftotext", "-layout", "-f", "1", "-l", "1", pdf, "-"],
                             check=True, capture_output=True, timeout=60)
        t = out.stdout.decode("utf-8", "replace")
    except Exception:
        return "", ""

    def grab(pattern):
        m = re.search(pattern, t, re.S | re.I | re.M)
        if not m:
            return ""
        return re.sub(r"\s+", " ", m.group(1)).strip()

    # Question runs until "Answer Statement:"; Words runs until the CCSS code or
    # the copyright line at the foot of the page.
    question = grab(r"Question:\s*(.*?)\s*Answer Statement:")
    # The footer sits after a run of blank lines: a CCSS code at column 0 and the
    # copyright. Cut at the blank-line gap rather than trying to match the code,
    # whose shape varies (6.NS.C.6a, 6.RP.A.3, 6.G.A.1 ...).
    words = grab(r"Words:\s*(.*?)(?=\n\s*\n|Original problems|©|\Z)")
    return question, words


def _prep_dare(url, workdir, tag):
    """DARE PDFs are 2 pages: problem, then answer. We use page 2 as the answer
    guide image, and lift the Question/Words text off page 1.
    Returns (answer_guide_png_or_None, question, words)."""
    raw = _download(url, os.path.join(workdir, f"{tag}.bin"))
    if not _is_pdf(raw):
        return None, "", ""
    pdf = os.path.join(workdir, f"{tag}.pdf")
    shutil.move(raw, pdf)
    question, words = _dare_text(pdf)
    pages = _pdf_pages_to_pngs(pdf, workdir, f"{tag}_p")
    if len(pages) >= 2:
        return BM.trim(pages[1]), question, words
    if pages:
        return BM.trim(pages[0]), question, words
    return None, question, words


def _prep_image(url, workdir, tag):
    """Math Talks and Sorts are stored as images, not PDFs."""
    raw = _download(url, os.path.join(workdir, f"{tag}.bin"))
    if _is_pdf(raw):
        pdf = os.path.join(workdir, f"{tag}.pdf")
        shutil.move(raw, pdf)
        pages = _pdf_pages_to_pngs(pdf, workdir, f"{tag}_p")
        return BM.trim(pages[0]) if pages else None
    png = os.path.join(workdir, f"{tag}.png")
    try:
        from PIL import Image
        Image.open(raw).convert("RGB").save(png)
    except Exception:
        return None
    return BM.trim(png)


# ---------- build ----------
def build_from_kit(kit, out_path, assets, dpi=150):
    """
    Build a lesson deck from one or more saved kits. Position is the lesson
    index: item i of each list becomes Lesson i.

    Preferred shape (three kits zipped by position):
      kit = {
        "title": str, "grade": int|str,
        "dares":      [{"name","standard","url"}, ...],
        "sorts":      [...],
        "math_talks": [...]
      }

    Legacy single-kit shape is still accepted:
      kit = {"title","grade","resource_type","items":[...]}

    Lists may differ in length. The deck runs for as many lessons as the
    LONGEST list; shorter lists simply leave that lesson's slides as empty
    dropzone placeholders rather than shifting items into the wrong lesson.
    """
    BM.ASSETS = assets

    dares = list(kit.get("dares") or [])
    sorts = list(kit.get("sorts") or [])
    talks = list(kit.get("math_talks") or [])

    # Legacy single-kit payload: route "items" into the right list.
    if not (dares or sorts or talks):
        rtype = (kit.get("resource_type") or "dare").strip()
        items = list(kit.get("items") or [])
        if rtype == "dare":
            dares = items
        elif rtype == "sort":
            sorts = items
        elif rtype == "math_talk":
            talks = items

    n_lessons = max(len(dares), len(sorts), len(talks))
    if n_lessons == 0:
        raise ValueError("Kit contains no items.")

    grade = kit.get("grade")
    try:
        grade = int(str(grade).strip())
    except Exception:
        pass

    # The teacher's own label for this kit ("1", "Unit 1", "Module 3"...). A kit
    # can span modules, so there is no module number to derive from the data --
    # whatever the teacher typed is what the deck shows. Falls back to the kit
    # title, then to a dash if neither is given.
    module_label = kit.get("module")
    module_label = str(module_label).strip() if module_label not in (None, "") else ""
    if not module_label:
        module_label = "—"

    tmp = tempfile.mkdtemp()
    img = os.path.join(tmp, "img")
    os.makedirs(img, exist_ok=True)

    def at(lst, i):
        return lst[i - 1] if i - 1 < len(lst) else None

    L_nums = list(range(1, n_lessons + 1))
    ccss_map = {}
    mt_by, sort_cells, ag_map = {}, {}, {}
    q_map, w_map = {}, {}

    for i in L_nums:
        d, s, t = at(dares, i), at(sorts, i), at(talks, i)

        # CCSS for the lesson footer: prefer the DARE's standard, then the
        # sort's, then the math talk's.
        std = ""
        for it in (d, s, t):
            if it and (it.get("standard") or "").strip():
                std = it["standard"].strip()
                break
        ccss_map[i] = std

        if d and d.get("url"):
            try:
                ag, q, w = _prep_dare(d["url"], img, f"d{i}")
                ag_map[i] = ag
                q_map[i], w_map[i] = q, w
            except Exception as e:
                print(f"  (lesson {i} dare: {e})")
        if t and t.get("url"):
            try:
                p = _prep_image(t["url"], img, f"t{i}")
                if p:
                    mt_by[i] = [p]
            except Exception as e:
                print(f"  (lesson {i} math talk: {e})")
        if s and s.get("url"):
            try:
                p = _prep_image(s["url"], img, f"s{i}")
                if p:
                    sort_cells[i] = BM.crop_sort_cells(p, img, f"s{i}")
            except Exception as e:
                print(f"  (lesson {i} sort: {e})")

    # Every lesson needs a math-talk slot so the plan stays uniform.
    for n in L_nums:
        mt_by.setdefault(n, [None])

    topics = BM.auto_topics(L_nums)

    plan = [(0, "toc")]
    for n in L_nums:
        secs = ["welcome", "mtr"] + [f"mtp{i}" for i in range(len(mt_by[n]))]
        if n in sort_cells:
            secs += ["sort", "rand"]
        secs += ["dareroutine", "dareguide", "dareedit", "game"]
        for sec in secs:
            plan.append((n, sec))

    idx = {(L, sec): i + 1 for i, (L, sec) in enumerate(plan)}
    welcome_idx = {n: idx[(n, "welcome")] for n in L_nums}
    codes = [c for c in ccss_map.values() if c]
    crange = f"{min(codes)} – {max(codes)}" if codes else ""

    MT = {
        "grade": grade,
        "module": module_label,
        "title": kit.get("title") or "Saved Kit",
        "topics": topics,
        "lesson_ccss": ccss_map,
        "welcome_idx": welcome_idx,
        "ccss_range": crange,
    }

    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)

    for (L, sec) in plan:
        if sec == "toc":
            BM.b_toc(prs, MT)
        elif sec == "welcome":
            has = L in sort_cells
            chip = {"Math Talk": idx[(L, "mtr")],
                    "DARE": idx[(L, "dareroutine")],
                    "Game": idx[(L, "game")]}
            if has:
                chip["Randomizer"] = idx[(L, "rand")]
                chip["Sort"] = idx[(L, "sort")]
            BM.b_welcome(prs, MT, L, has, chip)
        elif sec == "mtr":
            BM.b_mt_routine(prs, MT, L)
        elif sec.startswith("mtp"):
            BM.b_mt(prs, MT, L, mt_by[L][int(sec[3:])])
        elif sec == "sort":
            BM.b_sort(prs, MT, L, sort_cells.get(L))
        elif sec == "rand":
            BM.b_randroutine(prs, MT, L)
        elif sec == "dareroutine":
            BM.b_dare_routine(prs, MT, L, q_map.get(L, ""), w_map.get(L, ""))
        elif sec == "dareguide":
            BM.b_dareguide(prs, MT, L, ag_map.get(L))
        elif sec == "dareedit":
            BM.b_dareedit(prs, MT, L, ag_map.get(L))
        elif sec == "game":
            BM.b_game(prs, MT, L)

    raw = os.path.join(tmp, "raw.pptx")
    prs.save(raw)
    BM.postprocess(raw, out_path)
    return out_path
