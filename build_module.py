#!/usr/bin/env python3
"""
EM2 MODULE BUILDER — auto-loads a full interactive lesson deck from 4 source files.
See HANDOFF.md for usage. Run with --help for arguments.
"""
import os, re, sys, glob, shutil, subprocess, argparse, json, tempfile
import numpy as np
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image, ImageChops

HERE = os.path.dirname(os.path.abspath(__file__))

NAVY=RGBColor(0x18,0x30,0x5A); TEAL=RGBColor(0x2C,0x7A,0x7B)
RED=RGBColor(0xC0,0x39,0x2B); ORANGE=RGBColor(0xC9,0x82,0x1B); GOLD=RGBColor(0xB8,0x90,0x1A)
GREEN=RGBColor(0x2E,0x7D,0x32); BLUE=RGBColor(0x2F,0x6F,0xE0)
INK=RGBColor(0x27,0x31,0x3B); MUTED=RGBColor(0x7A,0x87,0x94); BODY=RGBColor(0x46,0x53,0x5F)
CARD=RGBColor(0xEE,0xF3,0xFA); LINE=RGBColor(0xD8,0xE0,0xEC); WHITE=RGBColor(0xFF,0xFF,0xFF)
CHEV=RGBColor(0xC2,0xCC,0xD8); TEAL_BAR=RGBColor(0x3F,0xB6,0xA8); CLOCK=RGBColor(0xCF,0xE0,0xFF)
SERIF="Georgia"; SANS="Calibri"
SW, SH = 10.0, 5.625
TOPIC_COLORS=[RED,ORANGE,GOLD,GREEN,BLUE]; TOPIC_NAMES=["TOPIC A","TOPIC B","TOPIC C","TOPIC D","TOPIC E"]
AG_CROP=(0.020,0.193,0.980,0.952)
PAGE_BOX=(0.35,1.10,9.30,3.95)   # full drawing area for a source page image

# ---------- value hygiene: never print None / "none" / "" ----------
def _toint(v):
    try: return int(str(v).strip())
    except Exception: return None

def _blank(v):
    return v is None or str(v).strip()=="" or str(v).strip().lower() in ("none","null","n/a","nan")

def norm_grade(v):
    if _blank(v): return None
    s=str(v).strip()
    if s.upper() in ("K","GK","KG"): return "K"
    m=re.search(r'(K|\d+)',s,re.I)
    if not m: return None
    g=m.group(1).upper()
    return "K" if g in ("K","0") else g

def norm_module(v):
    if _blank(v): return None
    m=re.search(r'\d+',str(v))
    return m.group(0) if m else None

def is_primary_grade(g):
    return str(g).strip().upper() in ("K","0","1","2")

def gml_label(grade,module,lesson=None,sep="   \u00b7   "):
    parts=[]
    if not _blank(grade): parts.append(f"Grade {grade}")
    if not _blank(module): parts.append(f"Module {module}")
    if lesson is not None: parts.append(f"Lesson {lesson}")
    return sep.join(parts)

def grade_from_filename(*paths):
    for p in paths:
        if not p: continue
        b=os.path.basename(str(p))
        m=(re.search(r'Grade\s*(K|\d+)',b,re.I) or re.search(r'\bG\s*(K|\d+)\b',b,re.I))
        if m: return norm_grade(m.group(1))
    return None

def module_from_filename(*paths):
    for p in paths:
        if not p: continue
        b=os.path.basename(str(p))
        m=(re.search(r'Module\s*(\d+)',b,re.I) or re.search(r'\bM\s*(\d+)\b',b,re.I))
        if m: return m.group(1)
    return None

CCSS_RE=re.compile(r'\b([K\d])\.([A-Z]{2,3})(?:\.([A-Z]))?\.(\d+[a-z]?)\b')
def ccss_of(text_):
    codes=[]
    for m in CCSS_RE.finditer(text_ or ""):
        code=".".join([x for x in m.groups() if x])
        if code not in codes: codes.append(code)
    return " \u2022 ".join(codes)

# ---------- source reading ----------
def shape_text(shapes):
    out=[]
    for sh in shapes:
        if sh.shape_type==6: out.append(shape_text(sh.shapes))
        if sh.has_table:
            for r in sh.table.rows:
                for c in r.cells: out.append(c.text)
        if sh.has_text_frame: out.append(sh.text_frame.text)
    return "\n".join(out)
def _is_pdf(path): return str(path).lower().endswith(".pdf")

def _pdf_texts(pdf):
    """Per-page text of a PDF via poppler's pdftotext."""
    out=[]
    n=0
    try:
        info=subprocess.run(["pdfinfo",pdf],check=True,capture_output=True,text=True).stdout
        m=re.search(r'Pages:\s*(\d+)',info)
        n=int(m.group(1)) if m else 0
    except Exception:
        n=0
    for p in range(1,n+1):
        try:
            t=subprocess.run(["pdftotext","-layout","-f",str(p),"-l",str(p),pdf,"-"],
                             check=True,capture_output=True,text=True).stdout
        except Exception:
            t=""
        out.append(t)
    return out

def slide_texts(src):
    if _is_pdf(src): return _pdf_texts(src)
    return [shape_text(s.shapes) for s in Presentation(src).slides]

def detect_dares(dares_src):
    """Returns (grade, module, lessons, texts). Each lesson carries its page index."""
    texts=slide_texts(dares_src); grade=module=None; lessons={}
    for i,t in enumerate(texts):
        if grade is None:
            gm=re.search(r'Grade\s*(K|\d+)',t,re.I)
            if gm: grade=norm_grade(gm.group(1))
        if module is None:
            mm=re.search(r'Module\s*(\d+)',t,re.I)
            if mm: module=mm.group(1)
        ln=re.search(r'Lesson\s*(\d+)',t,re.I)
        if not ln: continue
        n=int(ln.group(1))
        if n in lessons: continue
        q=re.search(r'L\d+\s*Question:\s*(.*?)\s*Answer Statement:',t,re.S)
        w=re.search(r'Words:\s*(.*?)\s*(?:CCSS|DARE\s*\u00a9|\Z)',t,re.S)
        lessons[n]={'n':n,'page':i,
                    'question':(q.group(1).strip() if q else ""),
                    'words':(re.sub(r'\s+',' ',w.group(1)).strip() if w else ""),
                    'ccss':ccss_of(t)}
    return grade, module, [lessons[k] for k in sorted(lessons)], texts

def detect_sort_lessons(sorts_pptx):
    texts=slide_texts(sorts_pptx); mapping={}; i=0
    while i<len(texts):
        lab=re.search(r'\bL(\d+)\b',texts[i])
        if lab and texts[i].count('L'+lab.group(1))>=3 and i+1<len(texts):
            mapping[int(lab.group(1))]=i+2; i+=2
        else: i+=1
    return mapping

# ---------- rendering / image ops ----------
RENDER_DPI = int(os.environ.get("RENDER_DPI", "120"))
JPEG_QUALITY = int(os.environ.get("JPEG_QUALITY", "80"))
GAME_RENDER_DPI = int(os.environ.get("GAME_RENDER_DPI", "90"))

# Progress hook. The API server sets build_module.PROGRESS to a callable so a
# background job can report honest, source-backed progress. Never a guessed %.
PROGRESS = None

def _progress(**kw):
    if PROGRESS:
        try: PROGRESS(**kw)
        except Exception: pass


def _pdf_of(src):
    """Returns a PDF path for src, converting a .pptx exactly once."""
    if _is_pdf(src): return src
    if src in _PDF_CACHE: return _PDF_CACHE[src]
    work=tempfile.mkdtemp(); shutil.copy(src,work)
    base=os.path.join(work,os.path.basename(src))
    subprocess.run(["soffice","--headless","--convert-to","pdf","--outdir",work,base],
                   check=True,capture_output=True,timeout=600,env={**os.environ,"HOME":work})
    pdf=base.rsplit('.',1)[0]+'.pdf'
    _PDF_CACHE[src]=pdf
    return pdf

_PDF_CACHE={}

def render_to_pngs(pptx,outdir,prefix,dpi=None):
    """Renders every page to JPEG (quality 80) at a modest dpi.

    JPEG at 120 dpi is 5-10x smaller than PNG at 150 dpi for scanned worksheet
    pages, with no visible loss at slide size. Each source PDF is opened once.
    """
    os.makedirs(outdir,exist_ok=True)
    pdf=_pdf_of(pptx)
    subprocess.run(["pdftoppm","-jpeg","-jpegopt",f"quality={JPEG_QUALITY}",
                    "-r",str(dpi or RENDER_DPI),pdf,os.path.join(outdir,prefix)],check=True)
    files=sorted(glob.glob(os.path.join(outdir,prefix+"*.jpg")))
    _progress(pages_rendered=len(files))
    return files

def trim(path,pad=8):
    with Image.open(path) as src:
        im=src.convert('RGB')
    bb=ImageChops.difference(im,Image.new('RGB',im.size,(255,255,255))).getbbox()
    if bb:
        l,t,r,b=bb; im=im.crop((max(0,l-pad),max(0,t-pad),min(im.width,r+pad),min(im.height,b+pad)))
    im.save(path,quality=JPEG_QUALITY,optimize=True)
    im.close()
    return path

def crop_box(path,frac,out):
    with Image.open(path) as src:
        im=src.convert('RGB')
    W,H=im.size; l,t,r,b=frac
    out=out.rsplit('.',1)[0]+'.jpg'
    cut=im.crop((int(W*l),int(H*t),int(W*r),int(H*b)))
    cut.save(out,quality=JPEG_QUALITY,optimize=True)
    cut.close(); im.close()
    return out

def _trim_im(im,pad=6):
    bb=ImageChops.difference(im.convert('RGB'),Image.new('RGB',im.size,(255,255,255))).getbbox()
    if bb:
        l,t,r,b=bb; im=im.crop((max(0,l-pad),max(0,t-pad),min(im.width,r+pad),min(im.height,b+pad)))
    return im

def crop_sort_cells(path,outdir,tag):
    """Cut the FIRST ROW of the 3-column sort sheet into 3 card images."""
    with Image.open(path) as src:
        im=src.convert('RGB')
    g=np.array(im.convert('L')); H,W=g.shape
    rowfrac=(g<180).mean(axis=1)
    cand=[y for y in range(int(H*0.04),int(H*0.6)) if rowfrac[y]>0.35]
    row_bottom=cand[0] if cand else H//7
    del g, rowfrac
    colw=W/3; cells=[]
    for c in range(3):
        x0=int(c*colw)+8; x1=int((c+1)*colw)-8
        cell=_trim_im(im.crop((x0,8,x1,row_bottom-4)))
        p=os.path.join(outdir,f'{tag}_cell{c}.jpg')
        cell.save(p,quality=JPEG_QUALITY,optimize=True); cell.close()
        cells.append(p)
    im.close()
    return cells

def detect_mathtalks(file,img,keep=None):
    """Group math-talk pages by lesson. Only pages for `keep` lessons are trimmed."""
    texts=slide_texts(file); pngs=render_to_pngs(file,img,'mt'); by={}
    for i,t in enumerate(texts):
        m=re.search(r'Lesson\s*(\d+)',t)
        if not m or i>=len(pngs): continue
        n=int(m.group(1))
        if keep is not None and n not in keep: continue
        by.setdefault(n,[]).append(trim(pngs[i]))
    return by,pngs

# ---------- drawing helpers ----------
def A(name): return os.path.join(ASSETS,name)
def _set_font(r,size,color,bold=False,italic=False,font=SANS):
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.name=font; r.font.color.rgb=color
def _sup(r): r.font._rPr.set('baseline','30000')
def fmt(sec): return f"{sec//60}:{sec%60:02d}"

def text(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=2):
    tb=s.shapes.add_textbox(I(l),I(t),I(w),I(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0); first=False
        for seg in line:
            r=p.add_run(); r.text=seg[0]
            _set_font(r,seg[1],seg[2],seg[3] if len(seg)>3 else False,seg[4] if len(seg)>4 else False,seg[5] if len(seg)>5 else SANS)
    return tb

def rect(s,l,t,w,h,fill=None,line_c=None,line_w=1.0,round_=False,shadow=False,radius=0.08):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,I(l),I(t),I(w),I(h))
    if round_:
        try: shp.adjustments[0]=radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line_c is None: shp.line.fill.background()
    else: shp.line.color.rgb=line_c; shp.line.width=Pt(line_w)
    shp.shadow.inherit=False
    if shadow:
        sp=shp._element.spPr; el=sp.makeelement(qn('a:effectLst'),{}); sp.append(el)
        sh=el.makeelement(qn('a:outerShdw'),{'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'}); el.append(sh)
        clr=sh.makeelement(qn('a:srgbClr'),{'val':'18305A'}); sh.append(clr); clr.append(clr.makeelement(qn('a:alpha'),{'val':'16000'}))
    return shp

def pill(s,l,t,w,h,fill,label,size=11,color=WHITE):
    shp=rect(s,l,t,w,h,fill=fill,round_=True,radius=0.5)
    tf=shp.text_frame; tf.word_wrap=False; tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=label; _set_font(r,size,color,bold=True); return shp

def pic(s,path,l,t,w=None,h=None):
    kw={}
    if w: kw['width']=I(w)
    if h: kw['height']=I(h)
    return s.shapes.add_picture(path,I(l),I(t),**kw)

def fit_pic(s,path,bl,bt,bw,bh,frame=False):
    with Image.open(path) as _probe: iw,ih=_probe.size
    ar=iw/ih; bar=bw/bh
    if ar>bar: w=bw; h=bw/ar
    else: h=bh; w=bh*ar
    l=bl+(bw-w)/2; t=bt+(bh-h)/2
    if frame: rect(s,l-0.06,t-0.06,w+0.12,h+0.12,fill=WHITE,line_c=LINE,line_w=1,round_=True,radius=0.03,shadow=True)
    return s.shapes.add_picture(path,I(l),I(t),width=I(w),height=I(h))

def logo(s): pic(s,A('logo.png'),SW-1.45,0.18,w=1.15)
def section_pill(s,label,fill): pill(s,0.42,0.30,1.4,0.36,fill,label,size=11)
def title(s,txt,color,size=30): text(s,1.0,0.30,8.0,0.6,[[(txt,size,color,True,False,SERIF)]],align=PP_ALIGN.CENTER)

# ---------- EDITABLE timer: navy pill + drawn clock + clickable time text + depleting bar ----------
def timer(s,seconds,pos='tr'):
    pw=1.45; ph=0.56
    pl=(SW-pw-1.55) if pos=='tr' else 0.42; pt=0.30
    rect(s,pl,pt,pw,ph,fill=NAVY,round_=True,radius=0.22)
    cd=0.26; cx=pl+0.15; cy=pt+(ph-cd)/2
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,I(cx),I(cy),I(cd),I(cd))
    o.fill.background(); o.line.color.rgb=CLOCK; o.line.width=Pt(1.5); o.shadow.inherit=False
    ccx=cx+cd/2; ccy=cy+cd/2
    for (x2,y2) in [(ccx,ccy-cd*0.34),(ccx+cd*0.27,ccy)]:
        h=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,I(ccx),I(ccy),I(x2),I(y2))
        h.line.color.rgb=CLOCK; h.line.width=Pt(1.25); h.shadow.inherit=False
    # editable time text (teacher can click and type a new value)
    tb=s.shapes.add_textbox(I(pl+0.40),I(pt),I(pw-0.46),I(ph)); tf=tb.text_frame
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.word_wrap=False
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=fmt(seconds); _set_font(r,22,WHITE,bold=True)
    tb.name="TimerText"
    bw=pw*0.86; bh=0.085; bl=pl+(pw-bw)/2; bt=pt+ph+0.05
    bar=rect(s,bl,bt,bw,bh,fill=TEAL_BAR,round_=True,radius=0.5); bar.name=f"TimerBar::{seconds}"
    return bar

def goto(shp,idx): shp.name=f"GOTO::{idx}"

def footer(s,ccss,grade,module,lesson):
    if not _blank(ccss):
        text(s,0.42,SH-0.42,4.5,0.3,[[("CCSS  ",9,INK,True),(str(ccss),9,MUTED)]])
    lab=gml_label(grade,module,lesson)
    if lab:
        text(s,SW-4.3,SH-0.42,3.88,0.3,[[(lab,9,MUTED)]],align=PP_ALIGN.RIGHT)

def routine_card(s,l,t,w,h,color,num,img,title_,desc):
    rect(s,l,t,w,h,fill=WHITE,line_c=LINE,line_w=1,round_=True,radius=0.06,shadow=True)
    c=rect(s,l+0.18,t+0.16,0.44,0.44,fill=color,round_=True,radius=0.5)
    tf=c.text_frame;tf.margin_top=0;tf.margin_bottom=0;pr=tf.paragraphs[0];pr.alignment=PP_ALIGN.CENTER
    rr=pr.add_run();rr.text=str(num);_set_font(rr,18,WHITE,bold=True)
    fit_pic(s,img,l+0.25,t+0.52,w-0.5,1.30)
    text(s,l+0.12,t+2.02,w-0.24,0.4,[[(title_,15.5,NAVY,True)]],align=PP_ALIGN.CENTER)
    text(s,l+0.2,t+2.46,w-0.4,0.95,[[(desc,12.5,BODY)]],align=PP_ALIGN.CENTER)

def reflected_callout(s,l,t,w,h):
    rect(s,l,t,w,h,fill=CARD,line_c=GREEN,line_w=2.0,round_=True,radius=0.14)
    b=s.shapes[-1]; tf=b.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="My edit was\u2026"; _set_font(r,18,GREEN,italic=True,font=SERIF)
    bx=l+0.45; by=t+h
    tri=s.shapes.build_freeform(I(bx),I(by),scale=I(1)/914400)
    tri.add_line_segments([(I(bx+0.30),I(by)),(I(bx),I(by+0.40)),(I(bx),I(by))],close=True)
    sh=tri.convert_to_shape(); sh.fill.solid(); sh.fill.fore_color.rgb=CARD
    sh.line.color.rgb=GREEN; sh.line.width=Pt(2.0); sh.shadow.inherit=False

def dropzone(s,l,t,w,h,label,field,fs=14):
    shp=rect(s,l,t,w,h,fill=RGBColor(0xF7,0xFA,0xFE),line_c=RGBColor(0xB7,0xC4,0xD6),line_w=1.75,round_=True,radius=0.04)
    ln=shp._element.spPr.find(qn('a:ln')); ln.append(ln.makeelement(qn('a:prstDash'),{'val':'dash'}))
    tf=shp.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=label; _set_font(r,fs,RGBColor(0x88,0x95,0xA6),bold=True)
    if field:
        p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
        r2=p2.add_run(); r2.text=field; _set_font(r2,10,RGBColor(0xA6,0xB2,0xC0),italic=True)

SCHED=[('Math Talk','sched_2.png',RED),('Randomizer','sched_3.png',ORANGE),
       ('Sort','sched_4.png',GOLD),('DARE','sched_5.png',GREEN),
       ('Task','work_team.png',TEAL),('Game','sched_6.png',BLUE)]

# ---------- slide builders ----------
def _toc_head(s,MT):
    logo(s)
    head=gml_label(MT['grade'],MT['module'],sep=" \u00b7 ")
    runs=[("EM",34,NAVY,True,False,SERIF),("2",20,NAVY,True,False,SERIF)]
    if head: runs.append((f"   {head}",34,NAVY,True,False,SERIF))
    text(s,0.6,0.60,8.8,0.8,[runs],align=PP_ALIGN.CENTER)
    for r in s.shapes[-1].text_frame.paragraphs[0].runs:
        if r.text=="2": _sup(r)
    text(s,0.6,1.36,8.8,0.4,[[(f"{MT['title']} \u2014 Table of Contents",15,TEAL,False,True,SERIF)]],align=PP_ALIGN.CENTER)

def _toc_foot(s,MT):
    if not _blank(MT['ccss_range']):
        text(s,0.42,SH-0.42,5,0.3,[[("CCSS  ",9,INK,True),(MT['ccss_range'],9,MUTED)]])
    lab=gml_label(MT['grade'],MT['module'])
    if lab: text(s,SW-4.0,SH-0.42,3.58,0.3,[[(lab,9,MUTED)]],align=PP_ALIGN.RIGHT)

def b_toc(prs,MT):
    s=prs.slides.add_slide(prs.slide_layouts[6]); _toc_head(s,MT)
    topics=MT['topics']
    if topics:
        n=len(topics); gap=0.16; mL=0.5; cw=(SW-2*mL-(n-1)*gap)/n; top=1.95; ch=2.95
        for i,tp in enumerate(topics):
            l=mL+i*(cw+gap); rect(s,l,top,cw,ch,fill=CARD,line_c=LINE,line_w=1,round_=True,radius=0.05)
            pill(s,l+0.12,top+0.14,cw-0.24,0.32,tp['color'],tp['name'],size=10.5); y=top+0.62
            for Ln in tp['lessons']:
                tb=text(s,l+0.1,y,cw-0.2,0.26,[[(f"Lesson {Ln}",11.5,NAVY)]],align=PP_ALIGN.CENTER)
                tb.text_frame.paragraphs[0].runs[0].font.underline=True; goto(tb,MT['welcome_idx'][Ln]); y+=0.275
    else:
        # No real topic data in the source: one plain lesson list, no invented groupings.
        L=list(MT['welcome_idx'].keys())
        rows=7; cols=max(1,-(-len(L)//rows)); rows=max(1,-(-len(L)//cols))
        mL=0.5; gap=0.14; cw=(SW-2*mL-(cols-1)*gap)/cols; top=1.95; rh=0.30
        for i,Ln in enumerate(L):
            c=i//rows; r=i%rows
            l=mL+c*(cw+gap); y=top+r*rh
            tb=text(s,l,y,cw,0.26,[[(f"Lesson {Ln}",12,NAVY)]],align=PP_ALIGN.CENTER)
            tb.text_frame.paragraphs[0].runs[0].font.underline=True; goto(tb,MT['welcome_idx'][Ln])
    _toc_foot(s,MT)

def b_welcome(prs,MT,L,parts,chip):
    """parts: ordered list of schedule step names actually in this lesson."""
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s)
    text(s,1.0,0.55,8.0,1.0,[[("Welcome!",54,NAVY,True,False,SERIF)]],align=PP_ALIGN.CENTER)
    lab=gml_label(MT['grade'],MT['module'],L,sep=" \u00b7 ")
    if lab: text(s,1.0,1.62,8.0,0.4,[[(lab,20,TEAL,False,True,SERIF)]],align=PP_ALIGN.CENTER)
    text(s,1.0,2.28,8.0,0.3,[[("TODAY'S SCHEDULE",12,MUTED,True)]],align=PP_ALIGN.CENTER)
    cards=[c for c in SCHED if c[0] in parts]
    cards.sort(key=lambda c: parts.index(c[0]))
    n=max(1,len(cards)); cw=1.45; gap=0.30; startl=(SW-(n*cw+(n-1)*gap))/2; top=2.7; ch=2.05
    for i,(name,icon,color) in enumerate(cards):
        l=startl+i*(cw+gap); rect(s,l,top,cw,ch,fill=WHITE,line_c=LINE,line_w=1,round_=True,radius=0.06,shadow=True)
        fit_pic(s,A(icon),l+0.1,top+0.12,cw-0.2,1.05)
        chs=pill(s,l+0.18,top+ch-0.5,cw-0.36,0.36,color,name,size=11)
        if name in chip: goto(chs,chip[name])
        if i<n-1: text(s,l+cw-0.02,top+0.7,gap+0.04,0.4,[[("\u203a",24,CHEV)]],align=PP_ALIGN.CENTER)
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_mt_routine(prs,MT,L):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"MATH TALK",RED)
    title(s,"Math Talk Routine",RED,size=25); timer(s,120)
    cards=[("mtr_1.png","Walk to the Math Talk","Gather and stand around the projected Math Talk."),
           ("mtr_2.png","Notice & Wonder","What do you notice? What do you wonder?"),
           ("mtr_3.png","Share, Listen, Connect","Share your thinking and build on others' ideas.")]
    cw=2.7; gap=0.3; startl=(SW-(3*cw+2*gap))/2; top=1.42; ch=3.6
    for i,(img,tt,d) in enumerate(cards): routine_card(s,startl+i*(cw+gap),top,cw,ch,RED,i+1,A(img),tt,d)
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_mt(prs,MT,L,img):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"MATH TALK",RED)
    title(s,"Math Talk",RED); timer(s,420)
    if img and os.path.exists(img): fit_pic(s,img,*PAGE_BOX,frame=True)
    else: dropzone(s,1.0,1.25,8.0,3.6,"Drop the Math Talk image here","mathTalkImg")
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_sort(prs,MT,L,cells):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); timer(s,120,pos='tl')
    text(s,1.8,0.32,6.5,0.5,[[("Do these cards match mathematically?",20,NAVY,True,False,SERIF)]],align=PP_ALIGN.CENTER)
    text(s,1.8,0.84,6.5,0.35,[[("Why or why not?",15,TEAL,False,True,SERIF)]],align=PP_ALIGN.CENTER)
    n=3; bw=2.5; gap=0.4; startl=(SW-(n*bw+(n-1)*gap))/2; top=1.5; bh=2.15
    for i in range(3):
        l=startl+i*(bw+gap)
        rect(s,l,top,bw,bh,fill=WHITE,line_c=BLUE,line_w=2.0,round_=True,radius=0.04)
        if cells and i<len(cells) and os.path.exists(cells[i]):
            fit_pic(s,cells[i],l+0.18,top+0.18,bw-0.36,bh-0.36)
        else:
            dropzone(s,l+0.12,top+0.12,bw-0.24,bh-0.24,f"Drop sort card {i+1}",f"sortCard{i+1}",fs=12)
    text(s,1.0,3.95,8.0,1.0,[[("Tell someone near you:",13,NAVY,True)],
        [("The cards match because…      The cards do ",12.5,BODY),("NOT",12.5,BODY,True),(" match because…",12.5,BODY)],
        [("Who wants to share their thinking?",12.5,TEAL,False,True)]],align=PP_ALIGN.CENTER,sp_after=4)
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_randroutine(prs,MT,L):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s)
    pill(s,0.42,0.30,1.6,0.36,ORANGE,"RANDOMIZER",size=10.5)
    text(s,2.1,0.32,0.4,0.34,[[("→",18,CHEV)]]); pill(s,2.5,0.30,1.0,0.36,GOLD,"SORT",size=10.5); timer(s,600)
    pw=4.25; ph=3.35; top=1.15
    rect(s,0.5,top,pw,ph,fill=CARD,line_c=LINE,line_w=1,round_=True,radius=0.05)
    text(s,0.75,top+0.14,pw-0.5,0.4,[[("Randomizer Routine",17,ORANGE,True,False,SERIF)]])
    text(s,0.82,top+0.6,pw-0.6,1.6,[[("1.  Choose a card from the basket.",12.5,BODY)],
        [("2.  Find two people whose cards mathematically match yours.",12.5,BODY)],
        [("3.  Bring the match to the teacher: \u201cOur cards match because\u2026\u201d",12.5,BODY)]],sp_after=6)
    for i,im in enumerate(["rand_1.png","rand_2.png","rand_3.png"]): fit_pic(s,A(im),0.7+i*1.22,top+2.15,1.12,1.0)
    rect(s,5.0,top,pw,ph,fill=CARD,line_c=LINE,line_w=1,round_=True,radius=0.05)
    text(s,5.25,top+0.14,pw-0.5,0.4,[[("Sort Routine",17,GOLD,True,False,SERIF)]])
    text(s,5.32,top+0.6,pw-0.6,1.5,[[("1.  With your group, prove your cards match and get the full sort from your teacher.",12.5,BODY)],
        [("2.  Work as a team to sort the cards in a way that makes sense to you.",12.5,BODY)]],sp_after=6)
    fit_pic(s,A('work_team.png'),6.55,top+2.05,1.6,1.15)
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_dare_worksheet(prs,MT,L,page):
    """Grades K\u20132: the DARE slide is the worksheet page itself, landscape, full area."""
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"DARE",GREEN)
    title(s,"DARE Routine",GREEN); timer(s,420)
    if page and os.path.exists(page): fit_pic(s,page,*PAGE_BOX,frame=True)
    else: dropzone(s,1.0,1.25,8.0,3.6,"Drop the DARE worksheet page here","dareWorksheetImg")
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_answer_page(prs,MT,L,page):
    """The lesson's answer-guide page, full area, after Time to Edit."""
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"DARE",GREEN)
    title(s,"DARE Answer Guide",GREEN,size=26); timer(s,180)
    if page and os.path.exists(page): fit_pic(s,page,*PAGE_BOX,frame=True)
    else: dropzone(s,1.0,1.25,8.0,3.6,"Drop the DARE answer-guide page here","dareAnswerPageImg")
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_task(prs,MT,L,page):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"TASK",TEAL)
    title(s,"Math Task",TEAL); timer(s,600)
    if page and os.path.exists(page): fit_pic(s,page,*PAGE_BOX,frame=True)
    else: dropzone(s,1.0,1.25,8.0,3.6,"Drop the Math Task page here","taskImg")
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_dare_routine(prs,MT,L,question,words):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"DARE",GREEN)
    title(s,"DARE Routine",GREEN); timer(s,420)
    fit_pic(s,A('dare_icons.png'),3.0,1.05,4.0,0.95)
    rect(s,0.7,2.1,8.6,2.95,fill=CARD,line_c=LINE,line_w=1,round_=True,radius=0.04)
    lines=[[("Question:  ",13.5,NAVY,True),(question or "",13.5,BODY)]]
    if not _blank(words):
        lines+=[[(" ",6,BODY)],[("Words:  ",13.5,NAVY,True),(words,13.5,BODY)]]
    text(s,1.0,2.32,8.0,2.55,lines,sp_after=4)
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_dareguide(prs,MT,L,ag):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"DARE",GREEN)
    title(s,"DARE Routine",GREEN); timer(s,360)
    fit_pic(s,A('dare_icons.png'),0.5,1.0,3.4,0.92)
    text(s,0.5,2.0,4.5,0.7,[[("Share & critique strategies.",16.5,NAVY,True,False,SERIF)],[("Look for your one edit!",16.5,NAVY,True,False,SERIF)]],sp_after=2)
    fit_pic(s,A('slc.png'),0.55,2.95,1.7,1.55)
    text(s,2.45,3.35,2.6,0.9,[[("What do you notice?",15.5,TEAL,True)],[("What do you wonder?",15.5,TEAL,True)]],sp_after=2)
    if ag and os.path.exists(ag): fit_pic(s,ag,5.13,0.95,4.55,4.05,frame=True)
    else: dropzone(s,5.13,0.95,4.55,4.05,"Drop the DARE answer-guide image here","dareAnswerImg")
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_dareedit(prs,MT,L,ag):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"DARE",GREEN)
    title(s,"DARE Routine",GREEN); timer(s,120)
    text(s,0.5,1.25,4.5,0.5,[[("Time to Edit!",26,NAVY,True,False,SERIF)]])
    text(s,0.5,1.85,4.5,0.7,[[("Use an edit pen or a different color to record new thinking, connections, or strategies.",13,BODY)]])
    fit_pic(s,A('edit_pens.png'),0.45,2.95,1.85,1.2)
    reflected_callout(s,2.45,3.0,2.6,0.92)
    if ag and os.path.exists(ag): fit_pic(s,ag,5.13,0.95,4.55,4.05,frame=True)
    else: dropzone(s,5.13,0.95,4.55,4.05,"Drop the DARE answer-guide image here","dareAnswerImg")
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def b_game(prs,MT,L):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"GAME",BLUE)
    title(s,"Game / Task Routine",BLUE,size=24); timer(s,420)
    cw=2.7; gap=0.3; startl=(SW-(3*cw+2*gap))/2; top=1.42; ch=3.6
    routine_card(s,startl,top,cw,ch,BLUE,1,A('gather.png'),"Gather Materials","Get your whiteboard, marker, eraser, and any tools you need.")
    routine_card(s,startl+(cw+gap),top,cw,ch,BLUE,2,A('work_team.png'),"Work with Your Team","Play the game or complete the task together, taking turns.")
    l3=startl+2*(cw+gap); rect(s,l3,top,cw,ch,fill=CARD,line_c=LINE,line_w=1,round_=True,radius=0.06)
    text(s,l3+0.2,top+0.35,cw-0.4,0.4,[[("Today's Game / Task",15,BLUE,True,False,SERIF)]],align=PP_ALIGN.CENTER)
    text(s,l3+0.25,top+1.05,cw-0.5,1.6,[[("Add today's specific game directions or task here.",13,RGBColor(0x88,0x95,0xA6),False,True)]],align=PP_ALIGN.CENTER)
    footer(s,MT['lesson_ccss'].get(L,""),MT['grade'],MT['module'],L)

def parse_game_links(raw):
    """One entry per line. A line that carries more than one URL is split so
    two pasted entries never render as a single run. PDF targets are dropped —
    those pages are already appended to the deck."""
    text=(raw or "").replace("\r\n","\n").replace("\r","\n")
    strip_chars=" |\t-\u2013\u2014:"
    entries=[]
    for line in text.split("\n"):
        line=line.strip()
        if not line: continue
        urls=list(re.finditer(r'https?://\S+',line))
        if len(urls)<=1:
            entries.append(line); continue
        start=0
        for m in urls:
            chunk=line[start:m.end()].strip(strip_chars)
            if chunk: entries.append(chunk)
            start=m.end()
        rest=line[start:].strip(strip_chars)
        if rest: entries.append(rest)
    out=[]; seen=set()
    for e in entries:
        m=re.search(r'https?://\S+',e)
        url=m.group(0).rstrip('.,;)') if m else ""
        label=(e.replace(m.group(0),"") if m else e).strip(strip_chars)
        label=re.sub(r'\s+',' ',label).strip() or url
        if url and re.search(r'\.pdf(\?|$)',url,re.I): continue
        key=(label.lower(),url)
        if key in seen: continue
        seen.add(key); out.append((label,url))
    return out

def b_game_divider(prs,MT):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s)
    text(s,0.8,2.05,8.4,0.9,[[("Game Directions & Videos",40,BLUE,True,False,SERIF)]],align=PP_ALIGN.CENTER)
    lab=gml_label(MT['grade'],MT['module'],sep=" \u00b7 ")
    if lab: text(s,0.8,3.05,8.4,0.4,[[(lab,16,TEAL,False,True,SERIF)]],align=PP_ALIGN.CENTER)

def b_game_links(prs,MT,links):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"GAMES",BLUE)
    title(s,"Game Directions & Videos",BLUE,size=26)
    tb=s.shapes.add_textbox(I(0.9),I(1.35),I(8.2),I(3.6)); tf=tb.text_frame
    tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    first=True
    for label,url in links:
        p=tf.paragraphs[0] if first else tf.add_paragraph()
        first=False; p.space_after=Pt(10)
        r=p.add_run(); r.text=label
        _set_font(r,16,BLUE,bold=True)
        if url:
            r.font.underline=True
            r.hyperlink.address=url

def b_game_page(prs,MT,page):
    s=prs.slides.add_slide(prs.slide_layouts[6]); logo(s); section_pill(s,"GAMES",BLUE)
    title(s,"Game Directions",BLUE,size=26)
    fit_pic(s,page,*PAGE_BOX,frame=True)

def auto_topics(L_nums):
    n=len(L_nums); k=min(5,n); size=-(-n//k); topics=[]
    for i in range(0,n,size):
        idx=len(topics)
        topics.append({'name':TOPIC_NAMES[idx] if idx<5 else f"TOPIC {idx+1}",'color':TOPIC_COLORS[idx%5],'lessons':L_nums[i:i+size]})
    return topics

def build(args):
    global ASSETS; ASSETS=args.assets
    g=lambda name,default=None: getattr(args,name,default)
    tmp=tempfile.mkdtemp(); img=os.path.join(tmp,'img'); os.makedirs(img,exist_ok=True)
    print("Reading DARE problems\u2026")
    grade,module,dares,dare_texts=detect_dares(args.dares)
    grade=norm_grade(g('grade')) or norm_grade(grade) or grade_from_filename(args.dares,args.mathtalks,g('answerguides'))
    module=norm_module(g('module')) or norm_module(module) or module_from_filename(args.dares,args.mathtalks,g('answerguides'))
    L_all=[d['n'] for d in dares]
    lo=_toint(g('lesson_from')); hi=_toint(g('lesson_to'))
    L_nums=[n for n in L_all if (lo is None or n>=lo) and (hi is None or n<=hi)]
    if not L_nums:
        raise ValueError(f"No lessons in the requested range ({lo or 'start'}-{hi or 'end'}). "
                         f"This module has lessons {min(L_all)}-{max(L_all)}.")
    keep=set(L_nums); N=len(L_nums)
    if len(L_nums)!=len(L_all):
        print(f"  Lesson range {lo or L_all[0]}-{hi or L_all[-1]}: building {N} of {len(L_all)} lessons")
    ccss_map={d['n']:d['ccss'] for d in dares}
    q_map={d['n']:d['question'] for d in dares}; w_map={d['n']:d['words'] for d in dares}
    page_map={d['n']:d['page'] for d in dares}
    print(f"  Grade {grade or '(unknown)'}, Module {module or '(unknown)'}, {N} lessons")
    _progress(lessons_total=N, stage="reading DARE problems")

    print("Rendering Math Talks\u2026"); _progress(stage="rendering math talks")
    mt_by,mt_png=detect_mathtalks(args.mathtalks,img,keep)
    if not mt_by:
        for i,n in enumerate(L_all):
            if n in keep: mt_by[n]=[trim(p) for p in mt_png[2*i:2*i+2]] or [None,None]
    for n in L_nums: mt_by.setdefault(n,[None,None])

    print("Rendering Sorts\u2026"); _progress(stage="rendering sorts")
    sort_cells={}
    if g('sorts'):
        try:
            sp=detect_sort_lessons(args.sorts); sorts=render_to_pngs(args.sorts,img,'so')
            for ln,front in sp.items():
                if ln in keep and 1<=front<=len(sorts): sort_cells[ln]=crop_sort_cells(trim(sorts[front-1]),img,f'L{ln}')
        except Exception as e: print("  (no sorts:",e,")")

    print("Rendering DARE pages\u2026"); _progress(stage="rendering DARE pages")
    dare_png=render_to_pngs(args.dares,img,'dr')
    dare_page={n:(trim(dare_png[page_map[n]]) if page_map[n]<len(dare_png) else None) for n in L_nums}

    print("Rendering DARE answer guides\u2026"); _progress(stage="rendering answer guides")
    ag_map={}
    same_file=(not g('answerguides')) or os.path.abspath(str(args.answerguides))==os.path.abspath(str(args.dares))
    if same_file:
        # One combined DARE Stories file: the page after a lesson page (no lesson
        # header of its own) is that lesson's answer page.
        lesson_pages=set(page_map.values())
        for n in L_nums:
            nxt=page_map[n]+1
            if nxt<len(dare_png) and nxt not in lesson_pages:
                ag_map[n]=trim(dare_png[nxt])
    else:
        ag=render_to_pngs(args.answerguides,img,'ag')
        for i,n in enumerate(L_all):
            if n in keep and i<len(ag): ag_map[n]=crop_box(ag[i],AG_CROP,os.path.join(img,f'agc_{n}.jpg'))

    task_page={}
    if g('tasks'):
        try:
            t_texts=slide_texts(args.tasks); t_png=render_to_pngs(args.tasks,img,'tk')
            for i,t in enumerate(t_texts):
                m=re.search(r'Lesson\s*(\d+)',t,re.I)
                if m and i<len(t_png) and int(m.group(1)) in keep:
                    task_page.setdefault(int(m.group(1)),trim(t_png[i]))
        except Exception as e: print("  (no tasks:",e,")")

    game_pages=[]
    if g('games'):
        try: game_pages=[trim(p) for p in render_to_pngs(args.games,img,'gm',dpi=GAME_RENDER_DPI)]
        except Exception as e: print("  (no game pages:",e,")")
    game_links=parse_game_links(g('game_links'))

    topics=None
    if g('topics') and os.path.exists(args.topics):
        tj=json.load(open(args.topics))
        topics=[{'name':t['name'],'color':TOPIC_COLORS[i%5],'lessons':t['lessons']} for i,t in enumerate(tj)]

    primary=is_primary_grade(grade) if grade else False

    # plan
    plan=[(0,'toc')]
    for n in L_nums:
        secs=['welcome','mtr']+[f'mtp{i}' for i in range(len(mt_by[n]))]
        if n in sort_cells: secs+=['sort','rand']
        secs+=['dareroutine','dareguide','dareedit']
        if ag_map.get(n): secs+=['answer']
        if n in task_page: secs+=['task']
        secs+=['game']
        for sec in secs: plan.append((n,sec))
    tail=[]
    if game_links or game_pages:
        tail.append((0,'gamedivider'))
        if game_links: tail.append((0,'gamelinks'))
        tail+= [(0,f'gamepage{i}') for i in range(len(game_pages))]
    plan+=tail
    idx={}
    for i,(L,sec) in enumerate(plan): idx[(L,sec)]=i+1
    welcome_idx={n:idx[(n,'welcome')] for n in L_nums}
    codes=[ccss_map[n] for n in L_nums if ccss_map.get(n)]
    crange=""
    if codes:
        flat=sorted({c.strip() for line in codes for c in line.split("\u2022")})
        crange=f"{flat[0]} \u2013 {flat[-1]}" if len(flat)>1 else flat[0]
    head=gml_label(grade,module,sep=" \u00b7 ")
    MT={'grade':grade,'module':module,'title':g('title') or head or "Daily Presentation Slides",
        'topics':topics,'lesson_ccss':ccss_map,'welcome_idx':welcome_idx,'ccss_range':crange}

    print("Building slides\u2026"); _progress(stage="building slides")
    prs=Presentation(); prs.slide_width=Emu(9144000); prs.slide_height=Emu(5143500)
    _done_lessons=set()
    for (L,sec) in plan:
        if L and L not in _done_lessons:
            _done_lessons.add(L); _progress(lessons_done=len(_done_lessons))
        if sec=='toc': b_toc(prs,MT)
        elif sec=='welcome':
            has=L in sort_cells
            parts=['Math Talk']
            chip={'Math Talk':idx[(L,'mtr')],'DARE':idx[(L,'dareroutine')],'Game':idx[(L,'game')]}
            if has:
                parts+=['Randomizer','Sort']
                chip['Randomizer']=idx[(L,'rand')]; chip['Sort']=idx[(L,'sort')]
            parts+=['DARE']
            if L in task_page:
                parts+=['Task']; chip['Task']=idx[(L,'task')]
            parts+=['Game']
            b_welcome(prs,MT,L,parts,chip)
        elif sec=='mtr': b_mt_routine(prs,MT,L)
        elif sec.startswith('mtp'): b_mt(prs,MT,L,mt_by[L][int(sec[3:])])
        elif sec=='sort': b_sort(prs,MT,L,sort_cells.get(L))
        elif sec=='rand': b_randroutine(prs,MT,L)
        elif sec=='dareroutine':
            if primary: b_dare_worksheet(prs,MT,L,dare_page.get(L))
            else: b_dare_routine(prs,MT,L,q_map.get(L,''),w_map.get(L,''))
        elif sec=='dareguide': b_dareguide(prs,MT,L,ag_map.get(L))
        elif sec=='dareedit': b_dareedit(prs,MT,L,ag_map.get(L))
        elif sec=='answer': b_answer_page(prs,MT,L,ag_map.get(L))
        elif sec=='task': b_task(prs,MT,L,task_page.get(L))
        elif sec=='game': b_game(prs,MT,L)
        elif sec=='gamedivider': b_game_divider(prs,MT)
        elif sec=='gamelinks': b_game_links(prs,MT,game_links)
        elif sec.startswith('gamepage'): b_game_page(prs,MT,game_pages[int(sec[8:])])
    raw=os.path.join(tmp,'raw.pptx'); prs.save(raw)
    _progress(stage="wiring timers & links", slides=len(plan))
    print("Wiring timers & links\u2026"); postprocess(raw,args.out); print("DONE \u2192",args.out)
    _progress(stage="done", slides=len(plan))

# ---------- post-process: timers + links ----------
TIMING='''<p:timing><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst><p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="5" presetID="22" presetClass="exit" presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="6" dur="1" fill="hold"><p:stCondLst><p:cond delay="{D}"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{S}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="hidden"/></p:to></p:set><p:animEffect transition="out" filter="wipe(right)"><p:cBhvr><p:cTn id="7" dur="{D}"/><p:tgtEl><p:spTgt spid="{S}"/></p:tgtEl></p:cBhvr></p:animEffect></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn><p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst><p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst><p:bldLst><p:bldP spid="{S}" grpId="0"/></p:bldLst></p:timing>'''
REL='<Relationship Id="{r}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="{t}"/>'

def postprocess(src,out):
    import zipfile
    work=tempfile.mkdtemp(); unp=os.path.join(work,'u'); os.makedirs(unp)
    with zipfile.ZipFile(src) as z: z.extractall(unp)
    pres=open(f'{unp}/ppt/presentation.xml').read(); rels=open(f'{unp}/ppt/_rels/presentation.xml.rels').read()
    r2t={m.group(1):m.group(2) for m in re.finditer(r'Id="(rId\d+)"[^>]*Target="slides/([^"]+)"',rels)}
    order=[r2t[m.group(1)] for m in re.finditer(r'<p:sldId[^>]*r:id="(rId\d+)"',pres)]
    pos2file={i+1:fn for i,fn in enumerate(order)}
    for pos,fn in pos2file.items():
        sp=f'{unp}/ppt/slides/{fn}'; xml=open(sp).read(); relp=f'{unp}/ppt/slides/_rels/{fn}.rels'
        rl=open(relp).read() if os.path.exists(relp) else '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"></Relationships>'
        mt=re.search(r'name="TimerBar::(\d+)"',xml)
        if mt:
            spid=re.search(r'<p:cNvPr id="(\d+)" name="TimerBar::\d+"',xml).group(1)
            xml=xml.replace('</p:sld>',TIMING.replace('{D}',str(int(mt.group(1))*1000)).replace('{S}',spid)+'</p:sld>',1)
        ex=[int(x) for x in re.findall(r'Id="rId(\d+)"',rl)]; nxt=[max(ex)+1 if ex else 1]; t2r={}; newr=[]
        def addlink(tp):
            tf=pos2file[tp]
            if tf in t2r: return t2r[tf]
            r=f'rId{nxt[0]}'; nxt[0]+=1; newr.append(REL.format(r=r,t=tf)); t2r[tf]=r; return r
        def repl(m):
            i,name=m.group(1),m.group(2); tp=int(name.split('::')[1])
            if tp==pos: return m.group(0)
            return f'<p:cNvPr id="{i}" name="{name}"><a:hlinkClick xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" r:id="{addlink(tp)}" action="ppaction://hlinksldjump"/></p:cNvPr>'
        xml=re.sub(r'<p:cNvPr id="(\d+)" name="(GOTO::\d+)"\s*/>',repl,xml)
        if newr: rl=rl.replace('</Relationships>',''.join(newr)+'</Relationships>'); open(relp,'w').write(rl)
        open(sp,'w').write(xml)
    if os.path.exists(out): os.remove(out)
    with zipfile.ZipFile(out,'w',zipfile.ZIP_DEFLATED) as z:
        for root,_,files in os.walk(unp):
            for f in files:
                full=os.path.join(root,f); z.write(full,os.path.relpath(full,unp))

if __name__=="__main__":
    ap=argparse.ArgumentParser()
    ap.add_argument('--mathtalks',required=True); ap.add_argument('--sorts',default=None)
    ap.add_argument('--dares',required=True); ap.add_argument('--answerguides',default=None)
    ap.add_argument('--tasks',default=None); ap.add_argument('--games',default=None)
    ap.add_argument('--game-links',dest='game_links',default=None)
    ap.add_argument('--out',required=True); ap.add_argument('--title',default=None)
    ap.add_argument('--topics',default=None); ap.add_argument('--grade',default=None)
    ap.add_argument('--lesson-from',dest='lesson_from',default=None)
    ap.add_argument('--lesson-to',dest='lesson_to',default=None)
    ap.add_argument('--module',default=None); ap.add_argument('--assets',default=os.path.join(HERE,'assets'))
    build(ap.parse_args())
